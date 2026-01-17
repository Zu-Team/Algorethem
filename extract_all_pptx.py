from pptx import Presentation
import os

def extract_pptx_content(pptx_path, output_file):
    """Extract text content from PowerPoint file and save to text file"""
    if not os.path.exists(pptx_path):
        print(f"Error: File not found at {pptx_path}")
        return False
    
    # Load the presentation
    prs = Presentation(pptx_path)
    
    # Extract text from all slides
    content = []
    content.append("=" * 80)
    content.append(output_file.replace('.txt', '').replace('_', ' ').title())
    content.append("=" * 80)
    content.append("")
    
    for i, slide in enumerate(prs.slides, 1):
        content.append(f"\n--- Slide {i} ---")
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    slide_text.append(text)
        content.append("\n".join(slide_text))
        content.append("")
    
    # Write to file
    with open(output_file, 'w', encoding='utf-8') as f:
        f.write("\n".join(content))
    
    print(f"Successfully extracted {len(prs.slides)} slides to {output_file}")
    return True

# Extract Chapter 5 - Greedy Algorithms
extract_pptx_content(
    r"ch5- Greedy Algorithms\ch5- Greedy Algorithms.pptx",
    "ch5_greedy_algorithms.txt"
)

# Extract Chapter 6 - Graphs
extract_pptx_content(
    r"ch6- graphs updated\ch6- graphs updated.pptx",
    "ch6_graphs.txt"
)

print("\nAll files extracted successfully!")

